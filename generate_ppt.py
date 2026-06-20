from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

try:
    prs = Presentation()

    slides_data = [
        {
            "title": "Billware - Smart Billing for Every Pharmacy",
            "content": "User Manual & Instructions",
            "screenshot": None
        },
        {
            "title": "1. Creating a New Store Account",
            "content": (
                "1. Open the Billware application.\n"
                "2. Click on the 'Create an account' link on the sign-in page.\n"
                "3. Enter your desired Account Username.\n"
                "4. Create a strong Password and re-enter it to confirm.\n"
                "5. Click the 'Create Account' button.\n\n"
                "Note: Each account operates as a completely independent, private store."
            ),
            "screenshot": "Screenshot 1: Create Account Page"
        },
        {
            "title": "2. Signing In to Your Dashboard",
            "content": (
                "1. Enter your Username and Password.\n"
                "2. Optional: Check 'Keep me signed in on this device' for faster access.\n"
                "3. Click the 'Sign In' button to access your dashboard.\n\n"
                "Tip: If you forgot your password, contact the pharmacy owner."
            ),
            "screenshot": "Screenshot 2: Sign In Page"
        },
        {
            "title": "3. Adding a New Customer",
            "content": (
                "1. Navigate to Home > Add Customer.\n"
                "2. Enter the customer's Name, Phone Number, and Address / Place. Mandatory fields are marked with a red asterisk (*).\n"
                "3. Click 'Save Customer' to add them to your database, or 'Cancel' to discard."
            ),
            "screenshot": "Screenshot 3: Add New Customer Page"
        },
        {
            "title": "4. Adding a New Doctor",
            "content": (
                "1. Navigate to Home > Add Doctor.\n"
                "2. Enter the Doctor Name, Phone Number, Clinic / Hospital Name, and Address / Place.\n"
                "3. Click 'Save Doctor' to register the doctor for future billing references."
            ),
            "screenshot": "Screenshot 4: Add New Doctor Page"
        },
        {
            "title": "5. Adding a Purchase (Part 1 - Bill Details)",
            "content": (
                "1. Navigate to Home > Purchase.\n"
                "2. Under Bill Details, enter the supplier's Bill No and the Firm (supplier name).\n"
                "3. Select the Date of purchase.\n"
                "4. Use 'Import Medicine List' to quickly load items, or clear them using 'Clear Imported'."
            ),
            "screenshot": "Screenshot 5: Add Purchase - Top"
        },
        {
            "title": "6. Adding a Purchase (Part 2 - Product Details)",
            "content": (
                "1. For each item, enter the Product Name, Batch No, Expire Date, Pack size, Quantity, MRP, Rate, and GST (%).\n"
                "2. The Total will be auto-calculated.\n"
                "3. Click '+ Add Another Item' if you have more products on the same bill.\n"
                "4. Review the Grand Total at the bottom and click 'Save All Purchases'."
            ),
            "screenshot": "Screenshot 6: Add Purchase - Bottom"
        },
        {
            "title": "7. Creating a New Sale Bill (Header Details)",
            "content": (
                "1. Navigate to Home > Sales. The Bill No and Date are auto-filled.\n"
                "2. Select an existing Customer Name from the dropdown, or click '+ New Customer' to add one.\n"
                "3. Select an existing Doctor Name, or click '+ New Doctor'."
            ),
            "screenshot": "Screenshot 7: New Sale Bill - Top"
        },
        {
            "title": "8. Creating a New Sale Bill (Medicine Details)",
            "content": (
                "1. In the Medicine Details section, select a product from the dropdown and choose the Batch No.\n"
                "2. The system auto-fills Exp Date, Pack, Stock, MRP, etc.\n"
                "3. Enter the Qty (quantity) you are selling.\n"
                "4. Click '+ Add Medicine Row' to add up to 15 items per bill.\n"
                "5. Check the Grand Total and click 'Save Bill' to finalize the sale."
            ),
            "screenshot": "Screenshot 8: New Sale Bill - Bottom"
        },
        {
            "title": "9. Viewing Reports & Analytics",
            "content": (
                "1. Navigate to Home > Report.\n"
                "2. Select a From Date and To Date, and choose the Type of transactions. Click 'Generate'.\n"
                "3. View summary cards: Total Sales, Total Purchase, Net Profit, and Total Items.\n"
                "4. Use Action Buttons to Print Report, Backup, or Restore data.\n"
                "5. Use 'Search Past Bills' to find historical bills."
            ),
            "screenshot": "Screenshot 9: Report Dashboard"
        },
        {
            "title": "10. Monitoring Medicine Stock",
            "content": (
                "1. On the Report page, click on 'Medicine Stock Report' to expand it.\n"
                "2. Use the search bar to find specific medicines by name or batch.\n"
                "3. Filter the view by clicking 'In Stock', 'Low Stock', 'Out of Stock', or 'Total Items'.\n"
                "4. Monitor the Total In, Total Out, and Current Stock columns to track inventory health."
            ),
            "screenshot": "Screenshot 10: Medicine Stock Report"
        },
        {
            "title": "11. Medicine Expiry Checker",
            "content": (
                "1. Navigate to the Expiry Checker tool.\n"
                "2. Select a timeframe to check (e.g., This Month, 1 Month, 2 Months, or Custom).\n"
                "3. Click 'Check Now' to generate the report.\n"
                "4. View the summary tabs and detailed list below to manage stock before it expires."
            ),
            "screenshot": "Screenshot 11: Medicine Expiry Checker"
        },
        {
            "title": "12. Transaction Records",
            "content": (
                "1. Navigate to Transaction Records to view a detailed log of all activities.\n"
                "2. Use the filters at the top to toggle between All, Sales, and Purchase records.\n"
                "3. View the summary statistics for Records, Quantity, Sales, Purchase, and Net amount.\n"
                "4. Use the red delete icon on the right to remove a specific record if needed."
            ),
            "screenshot": "Screenshot 12: Transaction Records"
        },
        {
            "title": "13. Inbox & Notifications Overview",
            "content": (
                "1. Click the Inbox icon in the top right to view alerts and notifications.\n"
                "2. The 'All' tab shows a chronological feed of all alerts.\n"
                "3. Use the 'Mark all read' button to clear the notifications.\n"
                "4. Notifications are color-coded (red for urgent, orange for warnings)."
            ),
            "screenshot": "Screenshot 13: Inbox panel - All"
        },
        {
            "title": "14. Managing Medicine Alerts",
            "content": (
                "1. Inside the Inbox, click the 'Medicines' tab to filter out business alerts.\n"
                "2. Quickly spot medicines that are Out of Stock or Expiring Soon.\n"
                "3. Use this view for daily inventory restocking and removing expiring items."
            ),
            "screenshot": "Screenshot 14: Inbox - Medicines tab"
        },
        {
            "title": "15. Tracking Business Updates",
            "content": (
                "1. Inside the Inbox, click the 'Business' tab to view operational alerts.\n"
                "2. Track milestones like total daily sales and performance metrics.\n"
                "3. Stay updated on store activity without leaving your current screen."
            ),
            "screenshot": "Screenshot 15: Inbox - Business tab"
        },
        {
            "title": "16. System Updates & Announcements",
            "content": (
                "1. Inside the Inbox, click the 'Updates' tab to view platform announcements.\n"
                "2. Read about new features, system maintenance, or important changes to Billware.\n"
                "3. Keep up to date with the latest tools designed to help your pharmacy."
            ),
            "screenshot": "Screenshot 16 & 17: Inbox - Updates"
        },
        {
            "title": "17. Managing Shop Settings",
            "content": (
                "1. Navigate to the Shop Settings page to manage your store identity.\n"
                "2. Review your Store Name, Invoice Prefix, Phone Number, and Address.\n"
                "3. Update Legal Information like Drug Licence No and GSTIN.\n"
                "4. Customize Invoice Terms & Conditions and Footer Message.\n"
                "5. Note: If editing is locked, click 'Request Edit Permission' to ask the Super Admin for access."
            ),
            "screenshot": "Screenshot 18: Shop Settings"
        }
    ]

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slides_data[0]["title"]
    subtitle.text = slides_data[0]["content"]

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[5] # Title only

    for slide_data in slides_data[1:]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = slide_data["title"]
        
        # Add a placeholder box for the screenshot
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(3.8)
        shape = shapes.add_shape(
            1, left, top, width, height # 1 is MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        text_frame = shape.text_frame
        text_frame.text = f"[ Please Replace This Box With {slide_data['screenshot']} ]"
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add text box for instructions
        left = Inches(0.5)
        top = Inches(5.2)
        width = Inches(9)
        height = Inches(2.0)
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = slide_data["content"]
        p.font.size = Pt(14)

    output_path = os.path.join(os.getcwd(), 'Billware_User_Instructions.pptx')
    prs.save(output_path)
    print(f"Success! Presentation saved to: {output_path}")

except Exception as e:
    print(f"Error generating PPT: {str(e)}")
